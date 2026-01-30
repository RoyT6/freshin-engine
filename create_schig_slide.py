#!/usr/bin/env python3
"""
Create SCHIG Architecture PowerPoint Slide (16:9)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(blank_layout)

# RGB color helper
def rgb_color(r, g, b):
    """Return RGB value for fill.fore_color.rgb"""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

# Colors
DARK_BLUE = rgb_color(0x1a, 0x36, 0x5d)
LIGHT_BLUE = rgb_color(0x4a, 0x90, 0xd9)
ACCENT_BLUE = rgb_color(0x5b, 0x9b, 0xd5)
DARK_GRAY = rgb_color(0x3d, 0x3d, 0x3d)
LIGHT_GRAY = rgb_color(0xf5, 0xf5, 0xf5)
WHITE = rgb_color(0xff, 0xff, 0xff)
GREEN = rgb_color(0x70, 0xad, 0x47)
ORANGE = rgb_color(0xed, 0x7d, 0x31)
PURPLE = rgb_color(0x70, 0x30, 0xa0)

def add_shape_with_text(slide, shape_type, left, top, width, height, text,
                        fill_color=None, font_size=12, font_bold=False,
                        font_color=DARK_GRAY, align=PP_ALIGN.CENTER):
    """Add a shape with centered text"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    shape.line.color.rgb = DARK_BLUE
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = font_bold
    p.font.color.rgb = font_color
    p.alignment = align

    # Vertical center
    tf.anchor = MSO_ANCHOR.MIDDLE

    return shape

def add_text_box(slide, left, top, width, height, text, font_size=11,
                 font_bold=False, font_color=DARK_GRAY, align=PP_ALIGN.LEFT):
    """Add a text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = font_bold
    p.font.color.rgb = font_color
    p.alignment = align
    return txBox

# === TITLE BANNER ===
title_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8)
)
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_BLUE
title_shape.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(8), Inches(0.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "SCHIG - Scared Child Gatherer"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

subtitle_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.25), Inches(4.5), Inches(0.4))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "Nightly Data Collection System | v1.0.0"
p.font.size = Pt(14)
p.font.color.rgb = rgb_color(0xbb, 0xbb, 0xbb)
p.alignment = PP_ALIGN.RIGHT

# === TASK SCHEDULER BOX (Top) ===
add_shape_with_text(
    slide, MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(4.5), Inches(1.0), Inches(4.3), Inches(0.5),
    "Windows Task Scheduler (02:00 AM)",
    fill_color=LIGHT_GRAY, font_size=10, font_bold=True
)

# Arrow down from scheduler
arrow1 = slide.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW, Inches(6.5), Inches(1.5), Inches(0.3), Inches(0.35)
)
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = DARK_BLUE
arrow1.line.fill.background()

# === MAIN SCHIG BOX ===
main_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(2.0), Inches(1.9), Inches(9.3), Inches(3.2)
)
main_box.fill.solid()
main_box.fill.fore_color.rgb = rgb_color(0xf0, 0xf4, 0xf8)
main_box.line.color.rgb = DARK_BLUE
main_box.line.width = Pt(2)

# schig.py label
add_text_box(slide, Inches(2.2), Inches(1.95), Inches(2), Inches(0.35),
             "schig.py", font_size=14, font_bold=True, font_color=DARK_BLUE)

# === ASYNC COLLECTORS SECTION ===
async_label = add_shape_with_text(
    slide, MSO_SHAPE.RECTANGLE,
    Inches(2.2), Inches(2.35), Inches(8.9), Inches(0.35),
    "ASYNC COLLECTORS (aiohttp + tenacity retry)",
    fill_color=ACCENT_BLUE, font_size=10, font_bold=True, font_color=WHITE
)

# Row 1 of collectors
collectors_row1 = [
    ("TMDB", GREEN),
    ("Streaming\nAvailability", GREEN),
    ("IMDB", LIGHT_BLUE),
    ("Rotten\nTomatoes", LIGHT_BLUE),
    ("HBO Max", LIGHT_BLUE),
]

x_start = 2.3
for i, (name, color) in enumerate(collectors_row1):
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_start + i * 1.75), Inches(2.8), Inches(1.6), Inches(0.6),
        name, fill_color=color, font_size=9, font_bold=True, font_color=WHITE
    )

# Row 2 of collectors
collectors_row2 = [
    ("Disney+", LIGHT_BLUE),
    ("Seeking\nAlpha", ORANGE),
    ("Twitter\nTrends", ORANGE),
    ("Sports\nHighlights", PURPLE),
    ("Reuters", PURPLE),
]

for i, (name, color) in enumerate(collectors_row2):
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_start + i * 1.75), Inches(3.5), Inches(1.6), Inches(0.6),
        name, fill_color=color, font_size=9, font_bold=True, font_color=WHITE
    )

# === SYNC COLLECTOR SECTION ===
sync_label = add_shape_with_text(
    slide, MSO_SHAPE.RECTANGLE,
    Inches(2.2), Inches(4.2), Inches(8.9), Inches(0.35),
    "SYNC COLLECTOR (HTTP Basic Auth)",
    fill_color=rgb_color(0x88, 0x88, 0x88), font_size=10, font_bold=True, font_color=WHITE
)

add_shape_with_text(
    slide, MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(5.3), Inches(4.6), Inches(2.7), Inches(0.4),
    "FlixPatrol", fill_color=GREEN, font_size=10, font_bold=True, font_color=WHITE
)

# === ARROW DOWN FROM MAIN BOX ===
arrow2 = slide.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW, Inches(6.5), Inches(5.15), Inches(0.3), Inches(0.4)
)
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = DARK_BLUE
arrow2.line.fill.background()

# === OUTPUT BOXES ===
outputs = [
    ("JSON\n(14 files)", Inches(3.0)),
    ("PARQUET\n(cranberry_fresh)", Inches(5.8)),
    ("CSV\n(2 files)", Inches(8.6)),
]

for text, left in outputs:
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left, Inches(5.65), Inches(2.0), Inches(0.7),
        text, fill_color=DARK_BLUE, font_size=10, font_bold=True, font_color=WHITE
    )

# === DESTINATION BOX ===
add_shape_with_text(
    slide, MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.45),
    "Fresh In!/  ->  FC-ALGO-80 Pipeline",
    fill_color=rgb_color(0x2d, 0x5a, 0x27), font_size=12, font_bold=True, font_color=WHITE
)

# === LEFT SIDE INFO PANEL ===
# API Sources
info_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.2), Inches(1.0), Inches(1.7), Inches(4.0)
)
info_box.fill.solid()
info_box.fill.fore_color.rgb = LIGHT_GRAY
info_box.line.color.rgb = DARK_BLUE

add_text_box(slide, Inches(0.3), Inches(1.05), Inches(1.5), Inches(0.3),
             "API SOURCES", font_size=10, font_bold=True, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)

api_list = """Priority 1:
  TMDB
  FlixPatrol
  Streaming Avail.

Priority 2:
  IMDB
  Rotten Tomatoes
  HBO Max
  Disney+

Priority 3-4:
  Seeking Alpha
  Twitter Trends
  Sports Highlights
  Reuters"""

add_text_box(slide, Inches(0.25), Inches(1.35), Inches(1.6), Inches(3.5),
             api_list, font_size=8, font_color=DARK_GRAY)

# === RIGHT SIDE INFO PANEL ===
# Countries
country_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(11.4), Inches(1.0), Inches(1.7), Inches(2.5)
)
country_box.fill.solid()
country_box.fill.fore_color.rgb = LIGHT_GRAY
country_box.line.color.rgb = DARK_BLUE

add_text_box(slide, Inches(11.5), Inches(1.05), Inches(1.5), Inches(0.3),
             "18 COUNTRIES", font_size=10, font_bold=True, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)

countries = """US  GB  CA  AU
DE  FR  IN  BR
JP  MX  ES  IT
NL  SE  NO  DK
PL  TR"""

add_text_box(slide, Inches(11.45), Inches(1.35), Inches(1.6), Inches(2.0),
             countries, font_size=9, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Schedule box
sched_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(11.4), Inches(3.6), Inches(1.7), Inches(1.4)
)
sched_box.fill.solid()
sched_box.fill.fore_color.rgb = LIGHT_GRAY
sched_box.line.color.rgb = DARK_BLUE

add_text_box(slide, Inches(11.5), Inches(3.65), Inches(1.5), Inches(0.3),
             "SCHEDULE", font_size=10, font_bold=True, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)

sched_text = """Daily @ 02:00 AM

Retry: 5 attempts
Backoff: 2-60s
Rate limits per API"""

add_text_box(slide, Inches(11.45), Inches(3.95), Inches(1.6), Inches(1.0),
             sched_text, font_size=8, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)

# === LEGEND ===
legend_y = Inches(5.15)
legend_items = [
    ("Priority 1", GREEN, Inches(0.3)),
    ("Priority 2", LIGHT_BLUE, Inches(1.5)),
    ("Priority 3", ORANGE, Inches(2.7)),
    ("Priority 4", PURPLE, Inches(3.9)),
]

for text, color, left in legend_items:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, legend_y, Inches(0.2), Inches(0.2))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    add_text_box(slide, left + Inches(0.25), legend_y - Inches(0.02), Inches(0.9), Inches(0.25),
                 text, font_size=8, font_color=DARK_GRAY)

# === FOOTER ===
footer = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.3))
tf = footer.text_frame
p = tf.paragraphs[0]
p.text = "FrameCore BFD | ViewerDBX Data Pipeline | Created: 2026-01-13"
p.font.size = Pt(9)
p.font.color.rgb = rgb_color(0x99, 0x99, 0x99)

# Save presentation
output_path = r"C:\Users\RoyT6\Downloads\Fresh In!\SCHIG_Architecture.pptx"
prs.save(output_path)
print(f"PowerPoint saved to: {output_path}")
